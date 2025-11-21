from pptx import Presentation
from pptx.util import Inches
import os
import re
import pandas as pd
import numpy

# Cleanup function DELETE LATER
if os.path.isdir("outputs/Presentación.pptx"):
	os.remove("outputs/Presentación.pptx")

image_filetypes = ["png", "jpg", "jpeg", "heic"]

prs = Presentation("inputs/Plantilla.pptx")

with open("config", 'r') as instruction_file:
	instructions = instruction_file.readlines()
n_instructions = len(instructions)

concursantes = pd.read_csv("inputs/csv/Concursantes.csv")
medallistas_individual = pd.read_csv("inputs/csv/Medallistas Individual.csv", header=0)
N_PRESIDIUM = 7

# location_match = re.compile("location: (.*)")
# location = location_match.match(instructions[0]).group(1)
# date_match = re.compile("date: (.*)")
# date = date_match.match(instructions[1]).group(1)

# DSL functions
# def add_location_and_date(slide):
# 	txBox = slide.shapes.add_textbox(height=Inches(0.4), left=Inches(0), top=Inches(7), width=Inches(9))
# 	tf = txBox.text_frame
# 	tf.text = f"{location}, {date}"

def add_title_slide():
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.placeholders[0].text = "CEREMONIA DE PREMIACIÓN"

def add_person_slide(i: int):
	slide = prs.slides.add_slide(prs.slide_layouts[8])
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		title_match = re.compile(r"\+ title: (.*)")
		title_get = title_match.match(instructions[i])
		if title_get: title = title_get.group(1)

		name_match = re.compile(r"\+ name: (.*)")
		name_get = name_match.match(instructions[i])
		if name_get: name = name_get.group(1)

		image_match = re.compile(r"\+ image: (.*)")
		image_get = image_match.match(instructions[i])
		if image_get: image = image_get.group(1)
		
		role_match = re.compile(r"\+ role: (.*)")
		role_get = role_match.match(instructions[i])
		if role_get: role = role_get.group(1)
		
		i+=1
		if i >= n_instructions: break
	
	if title: slide.placeholders[0].text = title
	if name: slide.placeholders[2].text = name
	try:
		if image: slide.placeholders[1].insert_picture(image)
	except: pass
	try:
		if role:
			add_role = slide.placeholders[2].text_frame.add_paragraph()
			add_role.text = role
			add_role.level = 1
	except: pass

def add_moment_slide(i: int):
	slide = prs.slides.add_slide(prs.slide_layouts[2])
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		name_match = re.compile(r"\+ name: (.*)")
		name_get = name_match.match(instructions[i])
		if name_get: name = name_get.group(1)

		# image_match = re.compile(r"\+ image: (.*)")
		# image_get = image_match.match(instructions[i])
		# if image_get: image = image_get.group(1)
		
		i+=1
		if i >= n_instructions: break
	
	if name: slide.placeholders[0].text = name
	# if image: slide.placeholders[1].insert_picture(image)

def add_parade():
	slide = prs.slides.add_slide(prs.slide_layouts[2])
	slide.placeholders[0].text = "DESFILE DE DELEGACIONES"
	estados_presentes = sorted(concursantes["Estado"].unique())
	for estado in estados_presentes:
		state_slide = prs.slides.add_slide(prs.slide_layouts[1])
		state_slide.placeholders[0].text = estado

		tf = state_slide.placeholders[1].text_frame
		is_first = True
		for concursante in concursantes[concursantes["Estado"] == estado]["NOMBRE COMPLETO"].values:
			if is_first:
				p = tf.paragraphs[0]
				is_first = False
			else:
				p = tf.add_paragraph()
			p.text = concursante

def add_individual_medals(i: int):
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		medal_match = re.compile(r"\+ medal: (.*)")
		medal_get = medal_match.match(instructions[i])
		if medal_get: medal = medal_get.group(1)

		i+=1
		if i >= n_instructions: break
	
	if not medal: return None

	section_slide = prs.slides.add_slide(prs.slide_layouts[2])
	section_slide.placeholders[0].text = f"Medallas de {medal}"

	medallistas = medallistas_individual[(medallistas_individual["Medalla"] == medal)].values
	blocks = [medallistas[i:i+N_PRESIDIUM] for i in range(0,len(medallistas), N_PRESIDIUM)]
	for block in blocks:
		for j in range(len(block)):
			medal_slide = prs.slides.add_slide(prs.slide_layouts[8])
			is_first = True
			for k in range(len(block)):
				medal_slide.placeholders[0].text = f"Medallas de {medal}".upper()
				tf = medal_slide.placeholders[2].text_frame
				if is_first:
					p = tf.paragraphs[0]
					is_first = False
				else:
					p = tf.add_paragraph()
				p.text = f"{block[k,3].upper()} ({block[k,0].upper()})"
				if j==k:
					p.font.bold = True
					for filetype in image_filetypes:
						try:
							picture_placeholder = medal_slide.placeholders[1].insert_picture(f"inputs/img/Individual/{block[k,2]}.{filetype}")
						except (FileNotFoundError, AttributeError):
							print(f"{block[k,2]}.{filetype} not found")
							pass

# Main DSL parsing loop
for i in range(n_instructions):
	if instructions[i][0] == "#": continue # Ignore comments
	if re.match("title", instructions[i]):
		add_title_slide()
	if re.match("parade", instructions[i]):
		add_parade()
	if re.match("person", instructions[i]):
		add_person_slide(i)
	if re.match("moment", instructions[i]):
		add_moment_slide(i)
	if re.match("individual", instructions[i]):
		add_individual_medals(i)
	

if not os.path.isdir("outputs"):
	os.mkdir("outputs")
prs.save("outputs/Presentación.pptx")
