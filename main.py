from pptx import Presentation
from pptx.util import Inches
import os
import re
import pandas as pd
import numpy

# Cleanup function DELETE LATER
if os.path.isdir("outputs/Presentación.pptx"):
	os.remove("outputs/Presentación.pptx")

prs = Presentation("inputs/Plantilla.pptx")

with open("config", 'r') as instruction_file:
	instructions = instruction_file.readlines()
n_instructions = len(instructions)

medallistas_individual = pd.read_csv("inputs/csv/Medallistas Individual.csv", header=0)
medallistas_equipos = pd.read_csv("inputs/csv/Medallistas Equipos.csv", header=0)
N_PRESIDIUM = 7

# location_match = re.compile("location: (.*)")
# location = location_match.match(instructions[0]).group(1)
# date_match = re.compile("date: (.*)")
# date = date_match.match(instructions[1]).group(1)

# DSL functions
# def add_location_and_date(slide):
# 	txBox = slide.shapes.add_textbox(height=Inches(0.4), left=Inches(0), top=Inches(7), width=Inches(9))
# 	tf = txBox.text_frame
# 	tf.text = f"{location}, {date}"

def add_title_slide():
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.placeholders[0].text = "CEREMONIA DE PREMIACIÓN"

def add_person_slide(i: int):
	slide = prs.slides.add_slide(prs.slide_layouts[8])
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		title_match = re.compile(r"\+ title: (.*)")
		title_get = title_match.match(instructions[i])
		if title_get: title = title_get.group(1)

		name_match = re.compile(r"\+ name: (.*)")
		name_get = name_match.match(instructions[i])
		if name_get: name = name_get.group(1)

		image_match = re.compile(r"\+ image: (.*)")
		image_get = image_match.match(instructions[i])
		if image_get: image = image_get.group(1)
		
		role_match = re.compile(r"\+ role: (.*)")
		role_get = role_match.match(instructions[i])
		if role_get: role = role_get.group(1)
		
		i+=1
		if i >= n_instructions: break
	
	if title: slide.placeholders[0].text = title
	if image: slide.placeholders[1].insert_picture(image)
	if name: slide.placeholders[2].text = name
	if role:
		add_role = slide.placeholders[2].text_frame.add_paragraph()
		add_role.text = role
		add_role.level = 1

def add_moment_slide(i: int):
	slide = prs.slides.add_slide(prs.slide_layouts[2])
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		name_match = re.compile(r"\+ name: (.*)")
		name_get = name_match.match(instructions[i])
		if name_get: name = name_get.group(1)

		# image_match = re.compile(r"\+ image: (.*)")
		# image_get = image_match.match(instructions[i])
		# if image_get: image = image_get.group(1)
		
		i+=1
		if i >= n_instructions: break
	
	if name: slide.placeholders[0].text = name
	# if image: slide.placeholders[1].insert_picture(image)

def add_individual_medals(i: int):
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		level_match = re.compile(r"\+ level: (.*)")
		level_get = level_match.match(instructions[i])
		if level_get: level = level_get.group(1)

		medal_match = re.compile(r"\+ medal: (.*)")
		medal_get = medal_match.match(instructions[i])
		if medal_get: medal = medal_get.group(1)

		i+=1
		if i >= n_instructions: break
	
	if not medal or not level: return None

	section_slide = prs.slides.add_slide(prs.slide_layouts[2])
	section_slide.placeholders[0].text = f"Medallas de {medal}\nNivel {level}"

	medallistas = medallistas_individual[(medallistas_individual["Medalla"] == medal) & (medallistas_individual["Nivel"] == level)].values
	blocks = [medallistas[i:i+N_PRESIDIUM] for i in range(0,len(medallistas), N_PRESIDIUM)]
	for block in blocks:
		for j in range(len(block)):
			medal_slide = prs.slides.add_slide(prs.slide_layouts[8])
			for k in range(len(block)):
				medal_slide.placeholders[0].text = f"Medallas de {medal}\nNivel {level}".upper()
				tf = medal_slide.placeholders[2].text_frame
				p = tf.add_paragraph()
				p.text = f"{block[k,3].upper()} ({block[k,0].upper()})"
				if j==k:
					p.font.bold = True
					picture_placeholder = medal_slide.placeholders[1].insert_picture(f"inputs/img/Individual/{block[k,2]}.png")


def add_team_medals(i: int):
	i += 1
	if i >= n_instructions: return None

	# Argument seeker
	while instructions[i][0] == "+":
		level_match = re.compile(r"\+ level: (.*)")
		level_get = level_match.match(instructions[i])
		if level_get: level = level_get.group(1)

		i+=1
		if i >= n_instructions: break
	
	if not level: return None
	
	section_slide = prs.slides.add_slide(prs.slide_layouts[2])
	section_slide.placeholders[0].text = f"Medallas por equipos\nNivel {level}"

	medallistas = medallistas_equipos[medallistas_equipos["Nivel"]==level]
	medallistas_bronce = medallistas_equipos[(medallistas_equipos["Nivel"]==level) & (medallistas_equipos["Medalla"] == "Bronce")].values
	for team in medallistas_bronce:
		slide = prs.slides.add_slide(prs.slide_layouts[8])
		slide.placeholders[0].text = "TERCER LUGAR"
		slide.placeholders[1].insert_picture(f"inputs/img/Teams/{team[0]}.png")
		slide.placeholders[2].text = team[1]

	medallistas_plata = medallistas_equipos[(medallistas_equipos["Nivel"]==level) & (medallistas_equipos["Medalla"] == "Plata")].values
	for team in medallistas_plata:
		slide = prs.slides.add_slide(prs.slide_layouts[8])
		slide.placeholders[0].text = "SEGUNDO LUGAR"
		slide.placeholders[1].insert_picture(f"inputs/img/Teams/{team[0]}.png")
		slide.placeholders[2].text = team[1]

	medallistas_oro = medallistas_equipos[(medallistas_equipos["Nivel"]==level) & (medallistas_equipos["Medalla"] == "Oro")].values
	for team in medallistas_oro:
		slide = prs.slides.add_slide(prs.slide_layouts[8])
		slide.placeholders[0].text = "PRIMER LUGAR"
		slide.placeholders[1].insert_picture(f"inputs/img/Teams/{team[0]}.png")
		slide.placeholders[2].text = team[1]
		

# Main DSL parsing loop
for i in range(n_instructions):
	if instructions[i][0] == "#": continue # Ignore comments
	if re.match("title", instructions[i]):
		add_title_slide()
	if re.match("person", instructions[i]):
		add_person_slide(i)
	if re.match("moment", instructions[i]):
		add_moment_slide(i)
	if re.match("individual", instructions[i]):
		add_individual_medals(i)
	if re.match("team", instructions[i]):
		add_team_medals(i)
	

if not os.path.isdir("outputs"):
	os.mkdir("outputs")
prs.save("outputs/Presentación.pptx")
